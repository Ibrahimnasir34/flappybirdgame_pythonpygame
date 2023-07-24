from pptx import Presentation

# Create a new PowerPoint presentation
presentation = Presentation()

# Slide 1: Title Slide
slide1 = presentation.slides.add_slide(presentation.slide_layouts[0])
title1 = slide1.shapes.title
subtitle1 = slide1.placeholders[1]

title1.text = "Blockchain Technology"
subtitle1.text = "An Introduction"

# Slide 2: What is Blockchain?
slide2 = presentation.slides.add_slide(presentation.slide_layouts[1])
title2 = slide2.shapes.title
content2 = slide2.placeholders[1]

title2.text = "What is Blockchain?"
content2.text = "Blockchain is a decentralized, distributed ledger technology that securely records and verifies transactions across multiple computers."

# Slide 3: Key Features of Blockchain
slide3 = presentation.slides.add_slide(presentation.slide_layouts[1])
title3 = slide3.shapes.title
content3 = slide3.placeholders[1]

title3.text = "Key Features of Blockchain"
content3.text = "1. Decentralization\n2. Transparency\n3. Immutability\n4. Security\n5. Smart Contracts"

# Slide 4: Use Cases of Blockchain
slide4 = presentation.slides.add_slide(presentation.slide_layouts[1])
title4 = slide4.shapes.title
content4 = slide4.placeholders[1]

title4.text = "Use Cases of Blockchain"
content4.text = "1. Cryptocurrencies\n2. Supply Chain Management\n3. Voting Systems\n4. Healthcare Records\n5. Intellectual Property Management"

# Slide 5: Benefits of Blockchain
slide5 = presentation.slides.add_slide(presentation.slide_layouts[1])
title5 = slide5.shapes.title
content5 = slide5.placeholders[1]

title5.text = "Benefits of Blockchain"
content5.text = "1. Enhanced Security\n2. Improved Efficiency\n3. Enhanced Transparency\n4. Reduced Costs\n5. Trusted Transactions"

# Save the PowerPoint presentation
presentation.save("blockchain_presentation.pptx")

# Display success message
print("Blockchain presentation created successfully!")
